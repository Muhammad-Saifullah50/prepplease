"""Deterministic document-extraction tools (contracts/tools.md, research R3).

Read-only (FR-019): these functions extract text from in-memory bytes and
never touch system state. The ingest pipeline runs them as preprocessing;
the parsing agent can also re-invoke them per page via registered tools.
"""

import asyncio
import io
from dataclasses import dataclass

import pypdfium2 as pdfium
from pptx import Presentation

from exambrain_agents.errors import ParsingFailedError

# A page whose extracted text is below this many characters is treated as
# image-only and routed to OCR (contracts/tools.md routing rule).
DIGITAL_TEXT_CHAR_THRESHOLD = 25


@dataclass(frozen=True)
class PageText:
    """Per-page extraction result."""

    page: int  # 1-based
    text: str
    char_count: int


@dataclass(frozen=True)
class SlideText:
    """Per-slide extraction result."""

    slide: int  # 1-based
    title: str | None
    text: str


def extract_pdf_text(pdf_bytes: bytes) -> list[PageText]:
    """pypdfium2 per-page text extraction (FR-001).

    Raises :class:`ParsingFailedError` on encrypted/corrupt/zero-page
    documents (edge case → paper marked failed by the pipeline).
    """
    try:
        doc = pdfium.PdfDocument(pdf_bytes)
    except Exception as exc:
        raise ParsingFailedError(
            f"unreadable PDF ({type(exc).__name__})"
        ) from exc
    try:
        if len(doc) == 0:
            raise ParsingFailedError("PDF has zero pages")
        pages: list[PageText] = []
        for i in range(len(doc)):
            try:
                text = doc[i].get_textpage().get_text_range()
            except Exception as exc:
                raise ParsingFailedError(
                    f"text extraction failed on page {i + 1} "
                    f"({type(exc).__name__})"
                ) from exc
            pages.append(PageText(page=i + 1, text=text, char_count=len(text)))
        return pages
    finally:
        doc.close()


def classify_pages(pages: list[PageText]) -> tuple[str, list[int]]:
    """Classify a PDF and list pages needing OCR (contracts/tools.md).

    Returns ``(document_type, ocr_page_numbers)`` — ``pdf_scanned`` when
    most pages fall below the digital-text character threshold.
    """
    ocr_pages = [
        p.page for p in pages if p.char_count < DIGITAL_TEXT_CHAR_THRESHOLD
    ]
    kind = "pdf_scanned" if len(ocr_pages) > len(pages) / 2 else "pdf_digital"
    return kind, ocr_pages


def _ocr_pages_sync(pdf_bytes: bytes, page_numbers: list[int]) -> list[PageText]:
    import pytesseract

    try:
        doc = pdfium.PdfDocument(pdf_bytes)
    except Exception as exc:
        raise ParsingFailedError(
            f"unreadable PDF ({type(exc).__name__})"
        ) from exc
    try:
        results: list[PageText] = []
        for number in page_numbers:
            bitmap = doc[number - 1].render(scale=2.0)  # ~144 dpi
            image = bitmap.to_pil()
            text = pytesseract.image_to_string(image)
            results.append(
                PageText(page=number, text=text, char_count=len(text))
            )
        return results
    finally:
        doc.close()


async def ocr_pdf_pages(
    pdf_bytes: bytes, page_numbers: list[int]
) -> list[PageText]:
    """Render pages with pypdfium2 and OCR with pytesseract (research R3).

    CPU-bound work runs in a thread so the event loop stays free
    (Constitution IV).
    """
    return await asyncio.to_thread(_ocr_pages_sync, pdf_bytes, page_numbers)


def extract_pptx_text(pptx_bytes: bytes) -> list[SlideText]:
    """python-pptx walk of slides → title + concatenated text frames."""
    try:
        prs = Presentation(io.BytesIO(pptx_bytes))
    except Exception as exc:
        raise ParsingFailedError(
            f"unreadable PPTX ({type(exc).__name__})"
        ) from exc
    slides: list[SlideText] = []
    for i, slide in enumerate(prs.slides, start=1):
        title: str | None = None
        texts: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text:
                continue
            if title is None and shape == slide.shapes.title:
                title = text
            else:
                texts.append(text)
        slides.append(SlideText(slide=i, title=title, text="\n".join(texts)))
    return slides
